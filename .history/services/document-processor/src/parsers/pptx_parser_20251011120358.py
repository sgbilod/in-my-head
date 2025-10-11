"""
PPTX Parser
Handles Microsoft PowerPoint .pptx files
"""

from pathlib import Path
import time

from .base_parser import BaseParser, ParsedDocument, ParsingError

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PptxParser(BaseParser):
    """Parser for Microsoft PowerPoint presentations"""
    
    @property
    def supported_formats(self) -> list[str]:
        return ['.pptx']
    
    async def parse(self, file_path: Path) -> ParsedDocument:
        """
        Parse PPTX file
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            ParsedDocument with extracted content
        """
        if not PPTX_AVAILABLE:
            raise ParsingError(
                "python-pptx library not installed",
                parser=self.name,
                file_path=str(file_path)
            )
        
        start_time = time.time()
        
        try:
            prs = Presentation(file_path)
            
            # Extract text from all slides
            text_parts = []
            tables = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    # Extract text from text frames
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        tables.append({
                            'slide': slide_num,
                            'data': table_data
                        })
                
                if slide_text:
                    text_parts.append(
                        f"[Slide {slide_num}]\n" + "\n".join(slide_text)
                    )
            
            text = "\n\n".join(text_parts)
            text = self._clean_text(text)
            
            # Extract metadata
            metadata = await self.extract_metadata(file_path)
            core_props = prs.core_properties
            
            metadata.update({
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': core_props.keywords or '',
                'created_at': core_props.created.isoformat()
                if core_props.created else '',
                'modified_at': core_props.modified.isoformat()
                if core_props.modified else '',
                'slide_count': len(prs.slides),
                'word_count': len(text.split()),
            })
            
            # Extract links
            links = self._extract_links(text)
            
            parsing_time = time.time() - start_time
            
            return ParsedDocument(
                text=text,
                text_length=len(text),
                metadata=metadata,
                tables=tables if tables else None,
                links=links if links else None,
                parser_used=self.name,
                parsing_time=parsing_time
            )
            
        except Exception as e:
            raise ParsingError(
                f"Failed to parse PPTX file: {str(e)}",
                parser=self.name,
                file_path=str(file_path)
            )
