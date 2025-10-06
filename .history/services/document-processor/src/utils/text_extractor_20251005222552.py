"""
Text Extraction Utility
Extracts text content from various document formats
"""

from pathlib import Path
from typing import Optional, Dict, Any
from datetime import datetime
import os


async def extract_text(file_path: str, doc_type: str) -> Optional[str]:
    """
    Extract text from a document
    
    Args:
        file_path: Path to document file
        doc_type: Document type (pdf, docx, txt, md, pptx, xlsx)
    
    Returns:
        Extracted text content
    """
    file_path = Path(file_path)
    
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    
    # Route to appropriate extractor
    if doc_type == 'txt' or doc_type == 'md':
        return await extract_text_file(file_path)
    elif doc_type == 'pdf':
        return await extract_pdf(file_path)
    elif doc_type == 'docx':
        return await extract_docx(file_path)
    elif doc_type == 'pptx':
        return await extract_pptx(file_path)
    elif doc_type == 'xlsx':
        return await extract_xlsx(file_path)
    else:
        # Try to read as text file
        try:
            return await extract_text_file(file_path)
        except:
            return None


async def extract_text_file(file_path: Path) -> str:
    """
    Extract text from plain text files
    
    Args:
        file_path: Path to text file
    
    Returns:
        File content as string
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(file_path, 'r', encoding='latin-1') as f:
            return f.read()


async def extract_pdf(file_path: Path) -> str:
    """
    Extract text from PDF files
    
    Args:
        file_path: Path to PDF file
    
    Returns:
        Extracted text content
    """
    try:
        import PyPDF2
        
        text = []
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text.append(page.extract_text())
        
        return '\n\n'.join(text)
    
    except ImportError:
        raise ImportError(
            "PyPDF2 is required for PDF extraction. "
            "Install it with: pip install PyPDF2"
        )
    except Exception as e:
        raise Exception(f"Failed to extract PDF: {str(e)}")


async def extract_docx(file_path: Path) -> str:
    """
    Extract text from DOCX files
    
    Args:
        file_path: Path to DOCX file
    
    Returns:
        Extracted text content
    """
    try:
        import docx
        
        doc = docx.Document(file_path)
        text = []
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text.append(' | '.join(row_text))
        
        return '\n\n'.join(text)
    
    except ImportError:
        raise ImportError(
            "python-docx is required for DOCX extraction. "
            "Install it with: pip install python-docx"
        )
    except Exception as e:
        raise Exception(f"Failed to extract DOCX: {str(e)}")


async def extract_pptx(file_path: Path) -> str:
    """
    Extract text from PowerPoint files
    
    Args:
        file_path: Path to PPTX file
    
    Returns:
        Extracted text content
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            slide_text.append(f"--- Slide {slide_num} ---")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(' | '.join(row_text))
            
            if len(slide_text) > 1:  # More than just the slide header
                text.append('\n'.join(slide_text))
        
        return '\n\n'.join(text)
    
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX extraction. "
            "Install it with: pip install python-pptx"
        )
    except Exception as e:
        raise Exception(f"Failed to extract PPTX: {str(e)}")


async def extract_xlsx(file_path: Path) -> str:
    """
    Extract text from Excel files
    
    Args:
        file_path: Path to XLSX file
    
    Returns:
        Extracted text content
    """
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = []
            sheet_text.append(f"--- Sheet: {sheet_name} ---")
            
            # Extract cell values
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None and str(cell).strip():
                        row_text.append(str(cell))
                
                if row_text:
                    sheet_text.append(' | '.join(row_text))
            
            if len(sheet_text) > 1:  # More than just the sheet header
                text.append('\n'.join(sheet_text))
        
        return '\n\n'.join(text)
    
    except ImportError:
        raise ImportError(
            "openpyxl is required for XLSX extraction. "
            "Install it with: pip install openpyxl"
        )
    except Exception as e:
        raise Exception(f"Failed to extract XLSX: {str(e)}")


def get_page_count(file_path: str, doc_type: str) -> Optional[int]:
    """
    Get page count for a document
    
    Args:
        file_path: Path to document
        doc_type: Document type
    
    Returns:
        Number of pages/slides/sheets or None
    """
    file_path = Path(file_path)
    
    if doc_type == 'pdf':
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                return len(pdf_reader.pages)
        except Exception:
            return None
    
    elif doc_type == 'pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            return len(prs.slides)
        except Exception:
            return None
    
    elif doc_type == 'xlsx':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            return len(wb.sheetnames)
        except Exception:
            return None
    
    # For other types, estimate from content
    return None


async def extract_metadata(file_path: str, doc_type: str) -> Dict[str, Any]:
    """
    Extract metadata from a document
    
    Args:
        file_path: Path to document
        doc_type: Document type
    
    Returns:
        Dictionary containing metadata
    """
    file_path = Path(file_path)
    metadata = {
        'file_size': os.path.getsize(file_path),
        'created': datetime.fromtimestamp(os.path.getctime(file_path)).isoformat(),
        'modified': datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat(),
    }
    
    if doc_type == 'pdf':
        metadata.update(await extract_pdf_metadata(file_path))
    elif doc_type == 'docx':
        metadata.update(await extract_docx_metadata(file_path))
    elif doc_type == 'pptx':
        metadata.update(await extract_pptx_metadata(file_path))
    elif doc_type == 'xlsx':
        metadata.update(await extract_xlsx_metadata(file_path))
    
    return metadata


async def extract_pdf_metadata(file_path: Path) -> Dict[str, Any]:
    """Extract metadata from PDF files"""
    try:
        import PyPDF2
        metadata = {}
        
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            
            if pdf_reader.metadata:
                if '/Author' in pdf_reader.metadata:
                    metadata['author'] = pdf_reader.metadata['/Author']
                if '/Title' in pdf_reader.metadata:
                    metadata['title'] = pdf_reader.metadata['/Title']
                if '/Subject' in pdf_reader.metadata:
                    metadata['subject'] = pdf_reader.metadata['/Subject']
                if '/CreationDate' in pdf_reader.metadata:
                    metadata['creation_date'] = str(pdf_reader.metadata['/CreationDate'])
                if '/ModDate' in pdf_reader.metadata:
                    metadata['mod_date'] = str(pdf_reader.metadata['/ModDate'])
        
        return metadata
    except Exception:
        return {}


async def extract_docx_metadata(file_path: Path) -> Dict[str, Any]:
    """Extract metadata from DOCX files"""
    try:
        import docx
        doc = docx.Document(file_path)
        metadata = {}
        
        if doc.core_properties:
            props = doc.core_properties
            if props.author:
                metadata['author'] = props.author
            if props.title:
                metadata['title'] = props.title
            if props.subject:
                metadata['subject'] = props.subject
            if props.created:
                metadata['created'] = props.created.isoformat()
            if props.modified:
                metadata['modified'] = props.modified.isoformat()
            if props.last_modified_by:
                metadata['last_modified_by'] = props.last_modified_by
        
        return metadata
    except Exception:
        return {}


async def extract_pptx_metadata(file_path: Path) -> Dict[str, Any]:
    """Extract metadata from PPTX files"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        metadata = {}
        
        if prs.core_properties:
            props = prs.core_properties
            if props.author:
                metadata['author'] = props.author
            if props.title:
                metadata['title'] = props.title
            if props.subject:
                metadata['subject'] = props.subject
            if props.created:
                metadata['created'] = props.created.isoformat()
            if props.modified:
                metadata['modified'] = props.modified.isoformat()
            if props.last_modified_by:
                metadata['last_modified_by'] = props.last_modified_by
        
        return metadata
    except Exception:
        return {}


async def extract_xlsx_metadata(file_path: Path) -> Dict[str, Any]:
    """Extract metadata from XLSX files"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        metadata = {}
        
        if wb.properties:
            props = wb.properties
            if props.creator:
                metadata['author'] = props.creator
            if props.title:
                metadata['title'] = props.title
            if props.subject:
                metadata['subject'] = props.subject
            if props.created:
                metadata['created'] = props.created.isoformat()
            if props.modified:
                metadata['modified'] = props.modified.isoformat()
            if props.lastModifiedBy:
                metadata['last_modified_by'] = props.lastModifiedBy
        
        return metadata
    except Exception:
        return {}
